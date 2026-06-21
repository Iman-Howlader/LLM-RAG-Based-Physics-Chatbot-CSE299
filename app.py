import hashlib
import logging
import mimetypes
import re
import secrets
import shutil
import sqlite3
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from fastapi import FastAPI, File, Form, Header, HTTPException, Query, Request, UploadFile
import os
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel

from docx import Document as DocxDocument
from langchain_chroma import Chroma
from langchain_community.document_loaders import PyPDFLoader
from langchain_core.documents import Document
from langchain_groq import ChatGroq
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
import pandas as pd


# Lightweight local mock for ChatGroq used during testing when a real API key
# is not available. Triggered when api_key == 'test_key' or env var MOCK_GROQ=1.
class MockChatGroq:
    def __init__(self, model_name: str = None, temperature: float = 0.1, groq_api_key: str = None):
        self.model_name = model_name
        self.temperature = temperature
        self.groq_api_key = groq_api_key

    def invoke(self, prompt: str):
        class _Resp:
            pass

        # Produce a short deterministic mocked answer referencing the prompt.
        resp = _Resp()
        resp.content = "Mocked answer: the uploaded documents mention limitations and performance issues." 
        return resp

logging.basicConfig(level=logging.INFO, format='[%(levelname)s] %(message)s')
logger = logging.getLogger(__name__)

app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)

CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200
MAX_HISTORY = 6
DEFAULT_TOP_K = 10
EMBEDDING_MODEL_NAME = "sentence-transformers/all-MiniLM-L6-v2"

STORAGE_ROOT = Path("storage")
UPLOADS_ROOT = STORAGE_ROOT / "uploads"
CHROMA_ROOT = STORAGE_ROOT / "chroma"
DB_PATH = STORAGE_ROOT / "multidoc.db"
USER_SESSIONS: Dict[str, Dict[str, Any]] = {}


def hash_password(password: str) -> str:
    return hashlib.sha256(password.encode("utf-8")).hexdigest()


def ensure_storage_dirs() -> None:
    STORAGE_ROOT.mkdir(parents=True, exist_ok=True)
    UPLOADS_ROOT.mkdir(parents=True, exist_ok=True)
    CHROMA_ROOT.mkdir(parents=True, exist_ok=True)


def get_db_connection() -> sqlite3.Connection:
    conn = sqlite3.connect(DB_PATH, check_same_thread=False)
    conn.row_factory = sqlite3.Row
    return conn


def initialize_database() -> None:
    ensure_storage_dirs()
    with get_db_connection() as conn:
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS users (
                email TEXT PRIMARY KEY,
                password_hash TEXT NOT NULL,
                token TEXT UNIQUE NOT NULL,
                api_key TEXT
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS documents (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                email TEXT NOT NULL,
                filename TEXT NOT NULL,
                file_path TEXT NOT NULL,
                total_chunks INTEGER NOT NULL,
                timestamp TEXT NOT NULL,
                UNIQUE(email, filename)
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS chat_history (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                email TEXT NOT NULL,
                role TEXT NOT NULL,
                message TEXT NOT NULL,
                timestamp TEXT NOT NULL
            )
            """
        )
        conn.commit()


def query_db(query: str, params: Tuple[Any, ...] = (), fetchone: bool = False, fetchall: bool = False, commit: bool = False):
    with get_db_connection() as conn:
        cursor = conn.execute(query, params)
        if commit:
            conn.commit()
        if fetchone:
            return cursor.fetchone()
        if fetchall:
            return cursor.fetchall()
        return None


def get_user_by_email(email: str) -> Optional[Dict[str, Any]]:
    row = query_db("SELECT * FROM users WHERE email = ?", (email,), fetchone=True)
    return dict(row) if row else None


def get_user_by_token(token: str) -> Optional[Dict[str, Any]]:
    row = query_db("SELECT * FROM users WHERE token = ?", (token,), fetchone=True)
    return dict(row) if row else None


def save_user(email: str, password_hash: str, token: str, api_key: Optional[str] = None) -> None:
    query_db(
        "INSERT INTO users (email, password_hash, token, api_key) VALUES (?, ?, ?, ?)",
        (email, password_hash, token, api_key),
        commit=True,
    )


def update_user_token(email: str, token: str) -> None:
    query_db(
        "UPDATE users SET token = ? WHERE email = ?",
        (token, email),
        commit=True,
    )


def update_user_api_key(email: str, api_key: str) -> None:
    query_db(
        "UPDATE users SET api_key = ? WHERE email = ?",
        (api_key, email),
        commit=True,
    )


def get_user_documents(email: str) -> List[Dict[str, Any]]:
    rows = query_db(
        "SELECT id, email, filename, file_path, total_chunks, timestamp FROM documents WHERE email = ? ORDER BY timestamp DESC",
        (email,),
        fetchall=True,
    )
    return [dict(row) for row in rows] if rows else []


def get_document(email: str, filename: str) -> Optional[Dict[str, Any]]:
    row = query_db(
        "SELECT * FROM documents WHERE email = ? AND filename = ?",
        (email, filename),
        fetchone=True,
    )
    return dict(row) if row else None


def insert_document(email: str, filename: str, file_path: str, total_chunks: int, timestamp: str) -> None:
    query_db(
        "INSERT INTO documents (email, filename, file_path, total_chunks, timestamp) VALUES (?, ?, ?, ?, ?)",
        (email, filename, file_path, total_chunks, timestamp),
        commit=True,
    )


def delete_document_record(email: str, filename: str) -> None:
    query_db(
        "DELETE FROM documents WHERE email = ? AND filename = ?",
        (email, filename),
        commit=True,
    )


def insert_chat_message(email: str, role: str, message: str, timestamp: str) -> None:
    query_db(
        "INSERT INTO chat_history (email, role, message, timestamp) VALUES (?, ?, ?, ?)",
        (email, role, message, timestamp),
        commit=True,
    )


def get_chat_history(email: str, limit: int = MAX_HISTORY) -> List[Dict[str, Any]]:
    rows = query_db(
        "SELECT role, message, timestamp FROM chat_history WHERE email = ? ORDER BY id DESC LIMIT ?",
        (email, limit),
        fetchall=True,
    )
    if not rows:
        return []
    history = [dict(row) for row in rows]
    return list(reversed(history))


def get_email_hash(email: str) -> str:
    return hashlib.sha256(email.encode("utf-8")).hexdigest()


def get_user_upload_dir(email: str) -> Path:
    directory = UPLOADS_ROOT / get_email_hash(email)
    directory.mkdir(parents=True, exist_ok=True)
    return directory


def get_user_chroma_dir(email: str) -> Path:
    directory = CHROMA_ROOT / get_email_hash(email)
    directory.mkdir(parents=True, exist_ok=True)
    return directory


def get_user_session(email: str) -> Dict[str, Any]:
    session = USER_SESSIONS.get(email)
    if not isinstance(session, dict):
        session = {
            "chat_history": [],
            "vectorstore": None,
            "vector_store": None,
            "uploaded_files": [],
        }
        USER_SESSIONS[email] = session

    if "vectorstore" not in session and "vector_store" in session:
        session["vectorstore"] = session["vector_store"]
    if "vector_store" not in session and "vectorstore" in session:
        session["vector_store"] = session["vectorstore"]

    session.setdefault("chat_history", [])
    session.setdefault("vectorstore", None)
    session.setdefault("vector_store", session["vectorstore"])
    session.setdefault("uploaded_files", [])
    return session


def load_user_vectorstore(email: str) -> Optional[Chroma]:
    chroma_dir = get_user_chroma_dir(email)
    try:
        if not chroma_dir.exists() or not any(chroma_dir.iterdir()):
            return Chroma(persist_directory=str(chroma_dir), embedding_function=embedding_model)
        return Chroma(persist_directory=str(chroma_dir), embedding_function=embedding_model)
    except Exception as exc:
        logger.warning("Unable to load or initialize Chroma vector store for %s: %s", email, exc)
        try:
            return Chroma(persist_directory=str(chroma_dir), embedding_function=embedding_model)
        except Exception as exc2:
            logger.error("Failed to create empty Chroma store for %s: %s", email, exc2)
            return None


def get_authenticated_email(authorization: str = Header(None), token: str = Query(None)) -> str:
    auth_token = None
    if authorization and authorization.startswith("Bearer "):
        auth_token = authorization.split(" ", 1)[1]
    elif token:
        auth_token = token

    if not auth_token:
        raise HTTPException(status_code=401, detail="Missing authentication token.")

    user = get_user_by_token(auth_token)
    if not user:
        raise HTTPException(status_code=401, detail="Invalid or expired token.")
    return user["email"]


@app.on_event("startup")
def startup() -> None:
    initialize_database()
    global embedding_model
    embedding_model = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL_NAME)


def clean_html(html: str) -> str:
    text = re.sub(r"<script.*?>.*?</script>", "", html, flags=re.S | re.I)
    text = re.sub(r"<style.*?>.*?</style>", "", text, flags=re.S | re.I)
    text = re.sub(r"<[^>]+>", " ", text)
    return re.sub(r"\s+", " ", text).strip()


def load_single_file(file_path: Path) -> List[Document]:
    ext = file_path.suffix.lower()
    docs: List[Document] = []
    try:
        if ext == ".pdf":
            docs = PyPDFLoader(str(file_path)).load()
        elif ext in (".txt", ".md"):
            text = file_path.read_text(encoding="utf-8", errors="ignore")
            docs = [Document(page_content=text)]
        elif ext in (".html", ".htm"):
            html = file_path.read_text(encoding="utf-8", errors="ignore")
            docs = [Document(page_content=clean_html(html))]
        elif ext == ".docx":
            d = DocxDocument(str(file_path))
            text = "\n".join(p.text for p in d.paragraphs if p.text)
            docs = [Document(page_content=text)]
        elif ext == ".csv":
            df = pd.read_csv(file_path, on_bad_lines="skip")
            docs = [Document(page_content=df.to_string(index=False))]
        elif ext == ".xlsx":
            df = pd.read_excel(file_path)
            docs = [Document(page_content=df.to_string(index=False))]
        elif ext == ".pptx":
            prs = Presentation(str(file_path))
            slides: List[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slides.append(shape.text)
            docs = [Document(page_content="\n\n".join(slides))]
        else:
            logger.warning("Unsupported file extension: %s", ext)
    except Exception as exc:
        logger.error("Error reading %s: %s", file_path.name, exc)

    valid_docs: List[Document] = []
    for doc in docs:
        content = str(doc.page_content or "").strip()
        if not content:
            continue
        doc.page_content = content
        doc.metadata["source"] = file_path.name
        valid_docs.append(doc)

    return valid_docs


def normalize_source(source: str) -> str:
    cleaned = re.sub(r"[^a-zA-Z0-9]+", "_", source.lower())
    return cleaned.strip("_") or "source"


def split_documents(documents: List[Document]) -> List[Document]:
    splitter = RecursiveCharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)
    all_chunks: List[Document] = []
    for doc in documents:
        if not doc.page_content.strip():
            continue
        try:
            chunks = splitter.split_documents([doc])
        except Exception as exc:
            logger.error("Chunking failed for source=%s error=%s", doc.metadata.get("source"), exc)
            continue
        for index, chunk in enumerate(chunks):
            text = str(chunk.page_content or "").strip()
            if not text:
                continue
            chunk.metadata = {
                **chunk.metadata,
                "source": doc.metadata.get("source", "unknown"),
                "page": doc.metadata.get("page", "unknown"),
                "chunk_id": f"{normalize_source(str(doc.metadata.get('source', 'unknown')))}_{doc.metadata.get('page', 'page')}_{index}",
            }
            chunk.page_content = text
            all_chunks.append(chunk)
    return all_chunks


def expand_query(question: str) -> List[str]:
    base = question.strip().rstrip("?.!")
    if not base:
        return [base]
    expansions = [base]
    expansions.extend([
        f"{base} summary",
        f"{base} key points",
        f"Details about {base}",
        f"{base} explanation",
    ])
    return list(dict.fromkeys(expansions))


def retrieve_relevant_chunks(vectorstore: Chroma, question: str, top_k: int) -> List[Tuple[Document, float]]:
    queries = expand_query(question)
    seen: Dict[Tuple[str, str], Tuple[float, Document]] = {}
    for query in queries:
        try:
            hits = vectorstore.similarity_search_with_score(query, k=top_k)
        except Exception as exc:
            logger.error("Similarity search failed for query=%s error=%s", query, exc)
            continue
        for doc, score in hits:
            if not getattr(doc, "page_content", "").strip():
                continue
            source = str(doc.metadata.get("source", "unknown"))
            chunk_id = str(doc.metadata.get("chunk_id", "unknown"))
            key = (source, chunk_id)
            if key not in seen or score < seen[key][0]:
                seen[key] = (score, doc)
    sorted_chunks = sorted(seen.values(), key=lambda item: item[0])[:top_k]
    ordered_chunks = [(doc, score) for score, doc in sorted_chunks]
    logger.info("Retrieved %s chunks for question=%s", len(ordered_chunks), question)
    return ordered_chunks


def build_prompt(question: str, chunks: List[Tuple[Document, float]], history: List[Dict[str, Any]]) -> str:
    context_parts = []
    for doc, _ in chunks:
        source = doc.metadata.get("source", "unknown")
        page = doc.metadata.get("page", "unknown")
        chunk_id = doc.metadata.get("chunk_id", "unknown")
        context_parts.append(f"[Source: {source}, Page: {page}, Chunk: {chunk_id}]\n{doc.page_content}")
    history_text = "\n".join(f"{item['role'].capitalize()}: {item['message']}" for item in history[-MAX_HISTORY:])
    return f"""You are a document-centric AI assistant. Answer using only the context below.
Rules:
- Use only the provided document context. Do not hallucinate.
- If the answer cannot be found, reply exactly: \"Sorry, I couldn't find an answer to that in your uploaded documents.\"
- Combine information across chunks and cite document names and pages.
- Keep the answer concise and factual.

Conversation history:
{history_text or 'None'}

Context:
{'\n\n---\n\n'.join(context_parts)}

Question: {question}
Answer:"""


@app.get("/")
def root() -> FileResponse:
    return FileResponse("index.html")


@app.post("/register")
async def register(email: str = Form(...), password: str = Form(...)) -> Dict[str, str]:
    email = email.strip().lower()
    if not email or not password:
        raise HTTPException(status_code=400, detail="Email and password are required.")
    if get_user_by_email(email):
        raise HTTPException(status_code=400, detail="An account with this email already exists.")
    token = secrets.token_urlsafe(32)
    save_user(email, hash_password(password), token)
    return {"message": "Account created successfully.", "token": token}


@app.post("/login")
async def login(email: str = Form(...), password: str = Form(...)) -> Dict[str, Optional[str]]:
    email = email.strip().lower()
    user = get_user_by_email(email)
    if not user or user["password_hash"] != hash_password(password):
        raise HTTPException(status_code=401, detail="Invalid email or password.")
    token = user.get("token")
    if not token:
        token = secrets.token_urlsafe(32)
        update_user_token(email, token)
    return {"token": token, "email": email, "api_key": user.get("api_key")}


@app.post("/logout")
def logout() -> Dict[str, str]:
    return {"message": "Logged out; clear the client token locally."}


@app.get("/me")
def me(authorization: str = Header(None), token: str = Query(None)) -> Dict[str, Any]:
    email = get_authenticated_email(authorization, token)
    user = get_user_by_email(email)
    return {"email": email, "api_key": user.get("api_key") if user else None}


@app.post("/upload")
async def upload_files(files: List[UploadFile] = File(...), authorization: str = Header(None), token: str = Query(None)) -> Dict[str, Any]:
    email = get_authenticated_email(authorization, token)
    user = get_user_by_email(email)
    if not user:
        raise HTTPException(status_code=401, detail="Invalid user.")

    user_upload_dir = get_user_upload_dir(email)
    user_chroma_dir = get_user_chroma_dir(email)
    existing_files = {doc["filename"] for doc in get_user_documents(email)}

    indexed_documents: List[Document] = []
    saved_docs: List[Dict[str, Any]] = []

    for upload_file in files:
        filename = Path(upload_file.filename.strip()).name
        if not filename:
            continue
        if filename in existing_files:
            logger.info("Skipping duplicate upload for %s: %s", email, filename)
            continue

        destination = user_upload_dir / filename
        with destination.open("wb") as buffer:
            shutil.copyfileobj(upload_file.file, buffer)

        docs = load_single_file(destination)
        if not docs:
            destination.unlink(missing_ok=True)
            continue

        chunked_docs = split_documents(docs)
        if not chunked_docs:
            destination.unlink(missing_ok=True)
            continue

        indexed_documents.extend(chunked_docs)
        insert_document(
            email,
            filename,
            str(destination),
            len(chunked_docs),
            datetime.utcnow().isoformat(),
        )
        saved_docs.append({
            "filename": filename,
            "file_path": str(destination),
            "total_chunks": len(chunked_docs),
        })

    if not indexed_documents:
        raise HTTPException(status_code=400, detail="No valid text was extracted from uploaded documents.")

    vectorstore = load_user_vectorstore(email)
    session = get_user_session(email)
    if vectorstore is None:
        vectorstore = Chroma.from_documents(
            indexed_documents,
            embedding_model,
            persist_directory=str(user_chroma_dir),
        )
        try:
            vectorstore.persist()
        except Exception as exc:
            logger.warning("Failed to persist new Chroma store for %s: %s", email, exc)
    else:
        vectorstore.add_documents(indexed_documents)
        try:
            vectorstore.persist()
        except Exception as exc:
            logger.warning("Failed to persist Chroma for %s: %s", email, exc)

    session["vectorstore"] = vectorstore
    session["vector_store"] = vectorstore
    session["uploaded_files"] = [doc["filename"] for doc in get_user_documents(email)]

    documents = get_user_documents(email)
    return {
        "message": "Upload complete.",
        "documents": documents,
        "files": [doc["filename"] for doc in documents],
    }


@app.get("/documents")
def list_documents(authorization: str = Header(None), token: str = Query(None)) -> Dict[str, List[Dict[str, Any]]]:
    email = get_authenticated_email(authorization, token)
    return {"documents": get_user_documents(email)}


@app.get("/download/{filename}")
def download_document(filename: str, inline: bool = Query(False), authorization: str = Header(None), token: str = Query(None)):
    """Download or view an uploaded document for the authenticated user.

    - `filename`: the stored filename (sanitized to prevent path traversal)
    - `inline`: when true, attempt to display inline; otherwise trigger attachment download
    """
    email = get_authenticated_email(authorization, token)
    filename = Path(filename).name
    document = get_document(email, filename)
    if not document:
        raise HTTPException(status_code=404, detail="Document not found.")

    file_path = Path(document["file_path"])
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File missing on server.")

    media_type = mimetypes.guess_type(str(file_path))[0] or "application/octet-stream"
    response = FileResponse(str(file_path), media_type=media_type)
    disposition = "inline" if inline else "attachment"
    response.headers["Content-Disposition"] = f'{disposition}; filename="{filename}"'
    return response


@app.get("/files/{filename}")
def view_file(filename: str, authorization: str = Header(None), token: str = Query(None)):
    """Compatibility endpoint used by the frontend to view files inline."""
    return download_document(filename, inline=True, authorization=authorization, token=token)


@app.delete("/delete/{filename}")
def delete_document(filename: str, authorization: str = Header(None), token: str = Query(None)) -> Dict[str, str]:
    email = get_authenticated_email(authorization, token)
    filename = Path(filename).name
    document = get_document(email, filename)
    if not document:
        raise HTTPException(status_code=404, detail="Document not found.")

    file_path = Path(document["file_path"])
    if file_path.exists():
        file_path.unlink()

    delete_document_record(email, filename)

    vectorstore = load_user_vectorstore(email)
    if vectorstore:
        try:
            vectorstore.delete(where={"source": filename})
            vectorstore.persist()
        except Exception as exc:
            logger.warning("Unable to delete chunks for %s: %s", filename, exc)
    session = USER_SESSIONS.get(email)
    if session is not None:
        session["vectorstore"] = None

    return {"message": f"Document '{filename}' deleted successfully."}


class ChatRequest(BaseModel):
    question: str
    api_key: Optional[str] = None
    model_name: str = "llama-3.3-70b-versatile"
    temperature: float = 0.1
    top_k: int = DEFAULT_TOP_K
    max_tokens: Optional[int] = None


@app.post("/chat")
async def chat(payload: ChatRequest, authorization: str = Header(None), token: str = Query(None)) -> Dict[str, Any]:
    email = get_authenticated_email(authorization, token)
    user = get_user_by_email(email)
    if not user:
        raise HTTPException(status_code=401, detail="Invalid user.")

    api_key = payload.api_key or user.get("api_key")
    if not api_key:
        raise HTTPException(status_code=400, detail="Groq API key is required.")
    if payload.api_key and payload.api_key != user.get("api_key"):
        update_user_api_key(email, payload.api_key)

    if email not in USER_SESSIONS:
        USER_SESSIONS[email] = {
            "chat_history": [],
            "vectorstore": None,
            "vector_store": None,
            "uploaded_files": [],
        }

    session = get_user_session(email)
    vectorstore = session.get("vectorstore") or session.get("vector_store")
    if vectorstore is None:
        vectorstore = load_user_vectorstore(email)
        session["vectorstore"] = vectorstore
        session["vector_store"] = vectorstore

    if vectorstore is None:
        raise HTTPException(status_code=400, detail="Please upload documents first.")

    if payload.top_k <= 0:
        payload.top_k = DEFAULT_TOP_K

    chat_history = get_chat_history(email, limit=MAX_HISTORY)
    retrieved_chunks = retrieve_relevant_chunks(vectorstore, payload.question, payload.top_k)
    if not retrieved_chunks:
        return {
            "answer": "Sorry, I couldn't find an answer to that in your uploaded documents.",
            "sources": [],
            "retrieved_chunks": [],
        }

    prompt = build_prompt(payload.question, retrieved_chunks, chat_history)
    logger.info("Chat request user=%s query=%s retrieved=%s prompt_len=%s", email, payload.question, len(retrieved_chunks), len(prompt))

    try:
        use_mock = (api_key == 'test_key') or os.getenv('MOCK_GROQ') == '1'
        if use_mock:
            llm = MockChatGroq(model_name=payload.model_name, temperature=payload.temperature, groq_api_key=api_key)
        else:
            llm = ChatGroq(model_name=payload.model_name, temperature=payload.temperature, groq_api_key=api_key)
        response = llm.invoke(prompt)
        answer_text = getattr(response, "content", None) or getattr(response, "text", None) or str(response)

        insert_chat_message(email, "user", payload.question, datetime.utcnow().isoformat())
        insert_chat_message(email, "assistant", answer_text, datetime.utcnow().isoformat())

        sources: List[str] = []
        retrieved_info: List[Dict[str, Any]] = []
        for doc, score in retrieved_chunks:
            source = str(doc.metadata.get("source", "unknown"))
            page = str(doc.metadata.get("page", "unknown"))
            chunk_id = str(doc.metadata.get("chunk_id", "unknown"))
            if source not in sources:
                sources.append(source)
            retrieved_info.append({
                "source": source,
                "page": page,
                "chunk_id": chunk_id,
                "score": score,
            })

        return {
            "answer": answer_text,
            "sources": sources,
            "retrieved_chunks": retrieved_info,
            "prompt_length": len(prompt),
        }
    except Exception as exc:
        logger.exception("Chat failed for user=%s", email)
        raise HTTPException(status_code=500, detail="An internal error occurred while processing your request.")


@app.exception_handler(HTTPException)
async def http_exception_handler(request: Request, exc: HTTPException):
    return JSONResponse(status_code=exc.status_code, content={"error": str(exc.detail)})


@app.exception_handler(Exception)
async def unhandled_exception_handler(request: Request, exc: Exception):
    logger.exception("Unhandled exception: %s", exc)
    return JSONResponse(status_code=500, content={"error": "An unexpected server error occurred."})
